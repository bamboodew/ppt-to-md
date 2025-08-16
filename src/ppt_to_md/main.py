#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT转Markdown工具

该脚本可以将PowerPoint文件(.pptx)转换为Markdown格式文件。
支持命令行参数和函数调用两种方式使用。
"""

import argparse
import sys
from typing import List
from pptx import Presentation


def ppt_to_md(ppt_path: str, md_path: str) -> bool:
    """
    将PPT文件转换为Markdown文件
    
    Args:
        ppt_path (str): PPT文件路径
        md_path (str): 输出的Markdown文件路径
        
    Returns:
        bool: 转换是否成功
    """
    try:
        prs = Presentation(ppt_path)
        md_content: List[str] = []
        
        for i, slide in enumerate(prs.slides, 1):
            md_content.append(f"## 幻灯片 {i}\n")  # 幻灯片标题
            
            # 收集幻灯片中的所有文本
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            
            # 如果幻灯片有内容，则添加到Markdown中
            if slide_texts:
                # 简单区分标题和内容（第一个非空文本作为标题）
                if len(slide_texts) > 1:
                    md_content.append(f"### {slide_texts[0]}\n")  # 标题
                    for text in slide_texts[1:]:
                        md_content.append(f"{text}\n")
                else:
                    md_content.append(f"{slide_texts[0]}\n")
                
                md_content.append("")  # 添加空行分隔
            else:
                md_content.append("_此幻灯片无内容_\n")
                md_content.append("")  # 添加空行分隔
        
        # 写入Markdown文件
        with open(md_path, "w", encoding="utf-8") as f:
            f.write("\n".join(md_content))
        
        print(f"成功转换 '{ppt_path}' 到 '{md_path}'")
        return True
        
    except Exception as e:
        print(f"转换失败: {str(e)}")
        return False


def main():
    """主函数，处理命令行参数"""
    parser = argparse.ArgumentParser(description="PPT转Markdown工具")
    parser.add_argument("input", help="输入的PPT文件路径")
    parser.add_argument("output", help="输出的Markdown文件路径")
    
    # 如果没有提供参数，则显示帮助信息
    if len(sys.argv) == 1:
        parser.print_help()
        return
    
    args = parser.parse_args()
    ppt_to_md(args.input, args.output)


if __name__ == "__main__":
    main()